"""
표준 문서 로더 및 청커
PPT, 엑셀 파일을 파싱하고 청킹하여 Vector Store에 인덱싱 가능한 형태로 변환
"""
from __future__ import annotations
from typing import List, Dict, Optional
from pathlib import Path
import json
import re
import sys

# 프로젝트 루트를 sys.path에 추가
project_root = Path(__file__).parent.parent.parent.parent.parent
sys.path.insert(0, str(project_root))

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False
    print("⚠️  pandas not installed. Excel parsing will be disabled.")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("⚠️  python-pptx not installed. PPT parsing will be disabled.")

try:
    # 최신 langchain에서는 langchain_core.documents에서 import
    try:
        from langchain_core.documents import Document
    except ImportError:
        # 구버전 호환성
        from langchain.schema import Document
    HAS_LANGCHAIN = True
except ImportError:
    HAS_LANGCHAIN = False
    print("⚠️  langchain not installed. Document creation will be disabled.")
    # 타입 힌트를 위한 더미 클래스
    class Document:
        pass

from src.project_generator.config import Config

try:
    from langchain_openai import ChatOpenAI
    HAS_LLM = True
except ImportError:
    HAS_LLM = False
    print("⚠️  langchain_openai not installed. LLM-based semantic text generation will be disabled.")


class StandardLoader:
    """
    표준 문서 로더
    PPT, 엑셀 파일을 파싱하고 청킹
    LLM을 사용하여 의미가 풍부한 semantic_text 생성
    """
    
    def __init__(self, enable_llm: bool = True):
        """
        초기화
        
        Args:
            enable_llm: LLM 기반 semantic_text 생성 활성화 여부
        """
        self.supported_formats = {
            '.xlsx': self._load_excel,
            '.xls': self._load_excel,
            '.pptx': self._load_ppt,
            '.txt': self._load_text,
            '.md': self._load_text,
            '.text': self._load_text,
        }
        
        # LLM 초기화 (semantic_text 생성용)
        self.enable_llm = enable_llm and HAS_LLM
        self.llm = None
        if self.enable_llm:
            try:
                self.llm = ChatOpenAI(
                    model=Config.DEFAULT_LLM_MODEL,
                    temperature=0.2,
                    top_p=1.0,
                    frequency_penalty=0.0,
                    presence_penalty=0.0
                )
            except Exception as e:
                print(f"⚠️  Failed to initialize LLM: {e}")
                self.enable_llm = False
    
    def load_standards(self, standards_path: Optional[Path] = None) -> List[Document]:
        """
        표준 문서 디렉토리에서 모든 파일을 로드하고 청킹
        
        Args:
            standards_path: 표준 문서 경로 (None이면 Config에서 가져옴)
            
        Returns:
            Document 리스트
        """
        if standards_path is None:
            standards_path = Config.COMPANY_STANDARDS_PATH
        
        if not standards_path.exists():
            print(f"⚠️  Standards path not found: {standards_path}")
            return []
        
        documents = []
        
        # 지원하는 형식의 파일 찾기
        for file_path in standards_path.rglob('*'):
            if file_path.is_file():
                # README 파일 제외 (README.md, README.txt, README 등)
                file_name_lower = file_path.name.lower()
                if file_name_lower.startswith('readme'):
                    continue
                
                suffix = file_path.suffix.lower()
                if suffix in self.supported_formats:
                    try:
                        loader_func = self.supported_formats[suffix]
                        chunks = loader_func(file_path)
                        documents.extend(chunks)
                        print(f"✅ Loaded {len(chunks)} chunks from {file_path.name}")
                    except Exception as e:
                        print(f"⚠️  Failed to load {file_path.name}: {e}")
        
        return documents
    
    def _load_excel(self, file_path: Path) -> List[Document]:
        """
        엑셀 파일을 파싱하고 청킹
        
        청킹 전략:
        1. 시트별로 분리
        2. 행 단위로 청킹 (각 행이 독립적인 표준 정보)
        """
        if not HAS_PANDAS:
            print("⚠️  pandas not installed. Cannot load Excel files.")
            return []
        
        if not HAS_LANGCHAIN:
            print("⚠️  langchain not installed. Cannot create Documents.")
            return []
        
        documents = []
        
        try:
            # 엑셀 파일 읽기 (모든 시트)
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # 빈 시트 스킵
                if df.empty:
                    continue
                
                # 청킹 전략: 항상 행 단위로 청킹 (각 행이 독립적인 표준 정보)
                # 이유: 섹션/카테고리별로 묶으면 하나의 거대한 chunk가 되어 유사도 검색이 부정확해짐
                # 예: "고객"을 검색해도 20개 도메인이 섞인 긴 텍스트와 비교하므로 유사도가 낮아짐
                chunks = self._chunk_excel_by_rows(df, file_path, sheet_name, chunk_size=1)
                documents.extend(chunks)
        
        except Exception as e:
            print(f"⚠️  Failed to parse Excel file {file_path}: {e}")
            import traceback
            traceback.print_exc()
        
        return documents
    
    def _chunk_excel_by_category(self, df: pd.DataFrame, file_path: Path, 
                                  sheet_name: str, category_col: str) -> List[Document]:
        """카테고리별로 청킹 (카테고리 사용하지 않음, 행 단위로 처리)"""
        # 카테고리 구분 없이 행 단위로 처리
        return self._chunk_excel_by_rows(df, file_path, sheet_name)
    
    def _chunk_excel_by_section(self, df: pd.DataFrame, file_path: Path,
                               sheet_name: str, section_col: str) -> List[Document]:
        """섹션별로 청킹"""
        documents = []
        
        for section, group in df.groupby(section_col):
            # 각 행을 자연어 텍스트와 구조화된 JSON으로 변환
            text_parts = []  # 자연어 텍스트 (임베딩용)
            structured_data_list = []  # JSON 데이터 (파싱용)
            
            for _, row in group.iterrows():
                text, structured_data = self._format_excel_row_as_standard_text(row, section)
                if text and text.strip():  # 자연어 텍스트가 있는 경우만
                    text_parts.append(text)
                    structured_data_list.append(structured_data)
            
            # 자연어 텍스트를 page_content에 저장 (임베딩용)
            # JSON은 metadata에만 저장 (파싱용)
            page_content_text = " ".join(text_parts) if text_parts else ""
            
            # JSON 배열을 metadata에 저장
            json_objects = []
            detected_domain = None  # 도메인 정보
            
            for structured_data in structured_data_list:
                json_obj = structured_data.copy()
                json_obj["section"] = str(section)
                json_objects.append(json_obj)
                
                # 도메인 정보 추출
                if not detected_domain and json_obj.get("domain"):
                    detected_domain = json_obj.get("domain")
                # entity_code에서 도메인 추출 (예: ODR, CPN, NTF)
                if not detected_domain and json_obj.get("entity_code"):
                    detected_domain = json_obj.get("entity_code")
            
            # 메타데이터 구성 (category 제거)
            metadata = {
                "type": "database_standard",  # 기본값
                "source": str(file_path),
                "sheet": sheet_name,
                "section": str(section),
                "row_count": len(group),
                "format": "excel",
                "structured_data": json.dumps(json_objects, ensure_ascii=False)  # JSON 형태로 메타데이터에 저장 (파싱용)
            }
            
            # 도메인 정보 추가
            if detected_domain:
                metadata["domain"] = detected_domain
            
            doc = Document(
                page_content=page_content_text,  # 자연어 텍스트 (임베딩용)
                metadata=metadata
            )
            documents.append(doc)
        
        return documents
    
    
    def _format_excel_row_as_standard_text(
        self, 
        row: pd.Series, 
        context: str = "",
        draft_context: Optional[Dict] = None
    ) -> tuple[str, Dict]:
        """
        엑셀 행을 구조화된 표준 텍스트와 JSON으로 변환
        
        Args:
            row: pandas Series (엑셀 행)
            context: 컨텍스트 (시트 이름 등, 참고용)
            draft_context: 초안 정보 (선택사항, 현재 사용하지 않음)
            
        Returns:
            (텍스트, 구조화된_데이터) 튜플
            - 텍스트: 임베딩 검색용 키워드 텍스트 (한글 제외, 영어만)
            - 구조화된_데이터: row의 모든 값을 그대로 저장한 딕셔너리
        """
        # LLM 사용하지 않음: row의 모든 값을 그대로 structured_data에 저장
        structured_data = {}
        # row의 모든 값을 structured_data에 저장 (컬럼명을 키로 사용)
        for col, val in row.items():
            if pd.notna(val) and str(val).strip():
                col_str = str(col).strip()
                val_str = str(val).strip()
                structured_data[col_str] = val_str
        
        # 검색용 텍스트 생성: row의 모든 값들을 공백으로 연결 (한글 포함)
        parts = []
        for col, val in row.items():
            if pd.notna(val) and str(val).strip():
                val_str = str(val).strip()
                parts.append(val_str)
        
        # 설명, 컨텍스트 등은 structured_data에만 저장
        # parts(검색용)에는 추가하지 않음 (노이즈 방지)
        if context:
            structured_data['context'] = context
        
        # type 설정 (기본값)
        if 'type' not in structured_data:
            structured_data['type'] = "database_standard"
        
        # ⚠️ 효율성 개선: 모든 원본 컬럼을 저장하지 않고 필요한 필드만 저장
        # LLM이 생성한 structured_data에는 이미 필요한 정보가 포함되어 있음
        # 원본 컬럼 데이터는 raw_data에만 저장 (필요시 참조용)
        # structured_data에는 LLM이 생성한 필드와 추출된 필드만 유지
        
        # 자연어 텍스트 생성 (임베딩용)
        # 간단한 키워드 텍스트만 사용 (LLM 불필요, 한글 포함)
        keyword_text = " ".join(parts) if parts else ""
        text = keyword_text
        
        # structured_data를 JSON 문자열로도 직렬화 (호환성을 위해)
        json_string = json.dumps(structured_data, ensure_ascii=False, indent=0)
        
        # (자연어 텍스트, 구조화된 데이터) 반환
        # 자연어 텍스트는 page_content에 저장되어 임베딩됨
        # structured_data는 metadata에 저장되어 LLM이 파싱함
        return text, structured_data
    
    def _chunk_excel_by_rows(self, df: pd.DataFrame, file_path: Path,
                            sheet_name: str, chunk_size: int = 20) -> List[Document]:
        """
        행 단위로 청킹 (기본 전략)
        chunk_size 행씩 묶어서 하나의 청크로
        """
        documents = []
        
            # 행을 chunk_size씩 묶어서 청킹
        for i in range(0, len(df), chunk_size):
            chunk_df = df.iloc[i:i+chunk_size]
            
            # 각 행을 자연어 텍스트와 구조화된 JSON으로 변환
            text_parts = []  # 자연어 텍스트 (임베딩용)
            structured_data_list = []  # JSON 데이터 (파싱용)
            # 시트 이름을 context로 전달 (참고용)
            context = sheet_name
            
            # 순차 처리: 각 행마다 LLM semantic_text 생성 (청킹 처리)
            for row_idx, row in chunk_df.iterrows():
                # LLM semantic_text 생성 (각 행마다 순차 처리)
                text, structured_data = self._format_excel_row_as_standard_text(row, context)
                if text and text.strip():  # 자연어 텍스트가 있는 경우만
                    text_parts.append(text)
                    structured_data_list.append(structured_data)
            
            # 자연어 텍스트를 page_content에 저장 (임베딩용)
            # JSON은 metadata에만 저장 (파싱용)
            page_content_text = " ".join(text_parts) if text_parts else ""
            
            # JSON 배열을 metadata에 저장
            json_objects = []
            detected_domain = None  # 도메인 정보
            
            for structured_data in structured_data_list:
                json_obj = structured_data.copy()
                json_objects.append(json_obj)
                
                # 도메인 정보 추출
                if not detected_domain and json_obj.get("domain"):
                    detected_domain = json_obj.get("domain")
                # entity_code에서 도메인 추출 (예: ODR, CPN, NTF)
                if not detected_domain and json_obj.get("entity_code"):
                    detected_domain = json_obj.get("entity_code")
            
            # 메타데이터 구성 (category 제거)
            metadata = {
                "type": "database_standard",  # 기본값
                "source": str(file_path),
                "sheet": sheet_name,
                "chunk_index": i // chunk_size,
                "row_start": i,
                "row_end": min(i + chunk_size, len(df)),
                "format": "excel",
                "structured_data": json.dumps(json_objects, ensure_ascii=False)  # JSON 형태로 메타데이터에 저장 (파싱용)
            }
            
            # 도메인 정보 추가
            if detected_domain:
                metadata["domain"] = detected_domain
            
            doc = Document(
                page_content=page_content_text,  # 자연어 텍스트 (임베딩용)
                metadata=metadata
            )
            documents.append(doc)
        
        return documents
    
    def _load_ppt(self, file_path: Path) -> List[Document]:
        """
        PPT 파일을 파싱하고 청킹
        
        청킹 전략:
        1. 슬라이드별로 분리
        2. 각 슬라이드의 텍스트를 추출
        3. 슬라이드 단위로 청킹
        """
        if not HAS_PPTX:
            print("⚠️  python-pptx not installed. Cannot load PPT files.")
            return []
        
        if not HAS_LANGCHAIN:
            print("⚠️  langchain not installed. Cannot create Documents.")
            return []
        
        documents = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_idx, slide in enumerate(prs.slides):
                # 슬라이드의 모든 텍스트 추출
                slide_texts = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_texts.append(text)
                
                # 텍스트가 있는 슬라이드만 처리
                if slide_texts:
                    slide_content = "\n".join(slide_texts)
                    
                    doc = Document(
                        page_content=slide_content,
                        metadata={
                            "type": "database_standard",
                            "source": str(file_path),
                            "slide_index": slide_idx,
                            "slide_count": len(prs.slides),
                            "format": "ppt"
                        }
                    )
                    documents.append(doc)
        
        except Exception as e:
            print(f"⚠️  Failed to parse PPT file {file_path}: {e}")
            import traceback
            traceback.print_exc()
        
        return documents
    
    def _load_text(self, file_path: Path) -> List[Document]:
        """
        텍스트 파일을 파싱하고 청킹
        
        청킹 전략:
        1. 섹션 단위로 분리 (##, ### 등 마크다운 헤더 또는 빈 줄)
        2. 문단 단위로 분리
        3. 규칙 단위로 분리 (번호 목록, 불릿 목록)
        """
        if not HAS_LANGCHAIN:
            print("⚠️  langchain not installed. Cannot create Documents.")
            return []
        
        documents = []
        
        try:
            # 파일 읽기 (여러 인코딩 시도)
            content = None
            encodings = ['utf-8', 'utf-8-sig', 'cp949', 'euc-kr', 'latin-1']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                print(f"⚠️  Failed to decode text file {file_path}")
                return []
            
            # 청킹 전략 1: 마크다운 헤더 기반 섹션 분리
            if self._has_markdown_headers(content):
                chunks = self._chunk_text_by_sections(content, file_path)
                documents.extend(chunks)
            
            # 청킹 전략 2: 규칙 단위 분리 (번호 목록, 불릿 목록)
            elif self._has_numbered_list(content) or self._has_bullet_list(content):
                chunks = self._chunk_text_by_rules(content, file_path)
                documents.extend(chunks)
            
            # 청킹 전략 3: 문단 단위 분리 (기본)
            else:
                chunks = self._chunk_text_by_paragraphs(content, file_path)
                documents.extend(chunks)
        
        except Exception as e:
            print(f"⚠️  Failed to parse text file {file_path}: {e}")
            import traceback
            traceback.print_exc()
        
        return documents
    
    def _has_markdown_headers(self, content: str) -> bool:
        """마크다운 헤더가 있는지 확인"""
        return bool(re.search(r'^#{1,6}\s+', content, re.MULTILINE))
    
    def _has_numbered_list(self, content: str) -> bool:
        """번호 목록이 있는지 확인"""
        return bool(re.search(r'^\d+[\.\)]\s+', content, re.MULTILINE))
    
    def _has_bullet_list(self, content: str) -> bool:
        """불릿 목록이 있는지 확인"""
        return bool(re.search(r'^[-*+]\s+', content, re.MULTILINE))
    
    def _chunk_text_by_sections(self, content: str, file_path: Path) -> List[Document]:
        """마크다운 헤더 기반 섹션별 청킹"""
        documents = []
        
        # 헤더로 섹션 분리
        # 패턴: ^#{1,6}\s+.*$ 다음 줄부터 다음 헤더까지
        sections = re.split(r'\n(?=#{1,6}\s+)', content)
        
        for i, section in enumerate(sections):
            section = section.strip()
            if len(section) < 50:  # 최소 길이
                continue
            
            # 첫 번째 줄이 헤더인지 확인
            lines = section.split('\n')
            header = lines[0] if lines else ""
            
            # 헤더 추출
            header_match = re.match(r'^(#{1,6})\s+(.+)$', header)
            if header_match:
                header_text = header_match.group(2)
            else:
                header_text = f"Section {i+1}"
            
            doc = Document(
                page_content=section,
                metadata={
                    "type": "database_standard",
                    "source": str(file_path),
                    "section": header_text,
                    "section_index": i,
                    "format": "text"
                }
            )
            documents.append(doc)
        
        return documents
    
    def _chunk_text_by_rules(self, content: str, file_path: Path) -> List[Document]:
        """규칙 단위로 청킹 (번호 목록, 불릿 목록)"""
        documents = []
        
        # 번호 목록 또는 불릿 목록으로 분리
        # 패턴: ^(\d+[\.\)]|[-*+])\s+ 로 시작하는 줄들
        lines = content.split('\n')
        current_rule = []
        rule_index = 0
        
        for line in lines:
            # 규칙 시작 (번호 또는 불릿)
            if re.match(r'^(\d+[\.\)]|[-*+])\s+', line):
                # 이전 규칙 저장
                if current_rule:
                    rule_text = '\n'.join(current_rule).strip()
                    if len(rule_text) > 30:  # 최소 길이
                        doc = Document(
                            page_content=rule_text,
                            metadata={
                                "type": "database_standard",
                                "source": str(file_path),
                                "rule_index": rule_index,
                                "format": "text"
                            }
                        )
                        documents.append(doc)
                        rule_index += 1
                
                # 새 규칙 시작
                current_rule = [line]
            else:
                # 규칙 내용 계속
                if current_rule:
                    current_rule.append(line)
        
        # 마지막 규칙 저장
        if current_rule:
            rule_text = '\n'.join(current_rule).strip()
            if len(rule_text) > 30:
                doc = Document(
                    page_content=rule_text,
                    metadata={
                        "type": "database_standard",
                        "source": str(file_path),
                        "rule_index": rule_index,
                        "format": "text"
                    }
                )
                documents.append(doc)
        
        return documents
    
    def _chunk_text_by_paragraphs(self, content: str, file_path: Path, 
                                   min_chunk_size: int = 200, 
                                   max_chunk_size: int = 1000) -> List[Document]:
        """
        문단 단위로 청킹 (기본 전략)
        
        Args:
            content: 텍스트 내용
            file_path: 파일 경로
            min_chunk_size: 최소 청크 크기 (문자 수)
            max_chunk_size: 최대 청크 크기 (문자 수)
        """
        documents = []
        
        # 빈 줄로 문단 분리
        paragraphs = content.split('\n\n')
        
        current_chunk = []
        current_size = 0
        chunk_index = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            para_size = len(para)
            
            # 현재 청크에 추가하면 최대 크기를 초과하는 경우
            if current_size + para_size > max_chunk_size and current_chunk:
                # 현재 청크 저장
                chunk_text = '\n\n'.join(current_chunk)
                if len(chunk_text) >= min_chunk_size:
                    doc = Document(
                        page_content=chunk_text,
                        metadata={
                            "type": "database_standard",
                            "source": str(file_path),
                            "chunk_index": chunk_index,
                            "format": "text"
                        }
                    )
                    documents.append(doc)
                    chunk_index += 1
                
                # 새 청크 시작
                current_chunk = [para]
                current_size = para_size
            else:
                # 현재 청크에 추가
                current_chunk.append(para)
                current_size += para_size
        
        # 마지막 청크 저장
        if current_chunk:
            chunk_text = '\n\n'.join(current_chunk)
            if len(chunk_text) >= min_chunk_size:
                doc = Document(
                    page_content=chunk_text,
                    metadata={
                        "type": "database_standard",
                        "source": str(file_path),
                        "chunk_index": chunk_index,
                        "format": "text"
                    }
                )
                documents.append(doc)
        
        return documents
    
    def determine_standard_type(self, file_path: Path) -> str:
        """
        파일명을 기반으로 표준 타입 결정
        
        Returns:
            "database_standard", "api_standard", "terminology_standard" 중 하나
        """
        file_name_lower = file_path.name.lower()
        
        if 'database' in file_name_lower or 'db' in file_name_lower or '테이블' in file_name_lower:
            return "database_standard"
        elif 'api' in file_name_lower or 'endpoint' in file_name_lower:
            return "api_standard"
        elif 'terminology' in file_name_lower or '용어' in file_name_lower or 'term' in file_name_lower:
            return "terminology_standard"
        else:
            # 기본값
            return "database_standard"

